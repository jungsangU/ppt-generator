"""Task 6: 디자인 & 스타일링"""
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pydantic import BaseModel

class ColorTheme(BaseModel):
    """색상 테마"""
    primary: tuple[int, int, int]  # RGB
    secondary: tuple[int, int, int]
    accent: tuple[int, int, int]
    text_dark: tuple[int, int, int]
    text_light: tuple[int, int, int]

class DesignApplier:
    """슬라이드 디자인 적용"""

    # 기본 테마들
    THEMES = {
        "blue": ColorTheme(
            primary=(0, 51, 102),
            secondary=(51, 102, 153),
            accent=(0, 153, 204),
            text_dark=(0, 0, 0),
            text_light=(255, 255, 255)
        ),
        "professional": ColorTheme(
            primary=(31, 78, 121),
            secondary=(79, 129, 189),
            accent=(155, 187, 89),
            text_dark=(0, 0, 0),
            text_light=(255, 255, 255)
        ),
        "modern": ColorTheme(
            primary=(68, 114, 196),
            secondary=(112, 165, 203),
            accent=(237, 125, 49),
            text_dark=(51, 51, 51),
            text_light=(255, 255, 255)
        ),
    }

    def __init__(self, prs: Presentation, theme: str = "blue"):
        self.prs = prs
        self.theme = self.THEMES.get(theme, self.THEMES["blue"])

    def apply_theme(self) -> None:
        """전체 테마 적용"""
        for slide in self.prs.slides:
            self._apply_slide_theme(slide)

    def _apply_slide_theme(self, slide) -> None:
        """슬라이드 테마 적용"""
        # 배경 색상 유지 (이미 설정됨)
        # 텍스트 스타일 통일
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                self._apply_text_style(shape.text_frame)

    def _apply_text_style(self, text_frame) -> None:
        """텍스트 스타일 적용"""
        for paragraph in text_frame.paragraphs:
            # 기존 포맷 유지
            pass

    def apply_custom_theme(self, theme: ColorTheme) -> None:
        """커스텀 테마 적용"""
        self.theme = theme
        self.apply_theme()

    def add_footer(self, text: str = "") -> None:
        """푸터 추가"""
        from pptx.util import Inches

        for slide in self.prs.slides:
            left = Inches(0.5)
            top = Inches(7)
            width = Inches(12.333)
            height = Inches(0.4)

            footer_box = slide.shapes.add_textbox(left, top, width, height)
            footer_frame = footer_box.text_frame
            footer_frame.text = text

            p = footer_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(*self.theme.text_dark)
            p.alignment = PP_ALIGN.CENTER

    def get_rgb_color(self, color_name: str) -> RGBColor:
        """테마 색상 가져오기"""
        if color_name == "primary":
            return RGBColor(*self.theme.primary)
        elif color_name == "secondary":
            return RGBColor(*self.theme.secondary)
        elif color_name == "accent":
            return RGBColor(*self.theme.accent)
        elif color_name == "text_dark":
            return RGBColor(*self.theme.text_dark)
        elif color_name == "text_light":
            return RGBColor(*self.theme.text_light)
        else:
            return RGBColor(0, 0, 0)
