"""Task 5: PPT 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

class PPTBuilder:
    """PowerPoint 프레젠테이션 생성"""

    def __init__(self, title: str = "프레젠테이션", width: float = 13.333, height: float = 7.5):
        self.prs = Presentation()
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)
        self.title = title

    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """타이틀 슬라이드 추가"""
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]  # 빈 레이아웃
        )

        # 배경
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 51, 102)

        # 제목
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(12.333)
        height = Inches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # 부제목
        if subtitle:
            left = Inches(0.5)
            top = Inches(4.2)
            width = Inches(12.333)
            height = Inches(1.0)
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle

            p = subtitle_frame.paragraphs[0]
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title: str, content: list[str], layout_type: str = "bullet") -> None:
        """콘텐츠 슬라이드 추가"""
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]  # 빈 레이아웃
        )

        # 배경
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # 제목
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.333)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)

        # 구분선
        line = slide.shapes.add_shape(
            1,  # 직선
            Inches(0.5), Inches(1.15),
            Inches(12.333), Inches(0)
        )
        line.line.color.rgb = RGBColor(0, 51, 102)
        line.line.width = Pt(2)

        # 콘텐츠
        if layout_type == "bullet":
            self._add_bullet_content(slide, content)
        elif layout_type == "two_column":
            self._add_two_column_content(slide, content)
        else:
            self._add_bullet_content(slide, content)

    def _add_bullet_content(self, slide, content: list[str]) -> None:
        """불릿 포인트 추가"""
        left = Inches(0.8)
        top = Inches(1.5)
        width = Inches(11.733)
        height = Inches(5.5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_before = Pt(6)
            p.space_after = Pt(6)

    def _add_two_column_content(self, slide, content: list[str]) -> None:
        """2단 콘텐츠 추가"""
        mid_point = len(content) // 2

        # 왼쪽
        left = Inches(0.8)
        top = Inches(1.5)
        width = Inches(5.5)
        height = Inches(5.5)

        left_box = slide.shapes.add_textbox(left, top, width, height)
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        for i, point in enumerate(content[:mid_point]):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 0, 0)

        # 오른쪽
        left = Inches(6.8)
        right_box = slide.shapes.add_textbox(left, top, width, height)
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        for i, point in enumerate(content[mid_point:]):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 0, 0)

    def add_closing_slide(self, text: str = "감사합니다") -> None:
        """마무리 슬라이드 추가"""
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        # 배경
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 51, 102)

        # 텍스트
        left = Inches(0.5)
        top = Inches(3)
        width = Inches(12.333)
        height = Inches(1.5)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = text

        p = text_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def save(self, output_path: str) -> None:
        """PPT 파일로 저장"""
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        print(f"PPT 저장 완료: {output_path}")
