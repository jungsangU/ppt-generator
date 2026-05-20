#!/usr/bin/env python3
"""PPT auto-generator: topic → Claude JSON → .pptx"""

import argparse
import json
import os
import sys

import anthropic
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

SYSTEM_PROMPT = """\
You are a professional presentation designer. Generate PowerPoint slide content as structured JSON.

The JSON must follow this exact schema:
{{
  "title": "Presentation title",
  "subtitle": "One-line tagline",
  "theme": {{
    "primary": "#hexcolor",
    "accent": "#hexcolor",
    "text_on_dark": "#hexcolor",
    "text_on_light": "#hexcolor"
  }},
  "slides": [
    {{ "type": "section", "title": "Section Name" }},
    {{
      "type": "content",
      "section": "Section Name",
      "title": "Slide Title",
      "layout": "bullets",
      "bullets": ["Point 1", "Point 2", "Point 3"],
      "sub_bullets": {{"Point 1": ["Sub A", "Sub B"]}},
      "highlight": "Optional key stat or callout"
    }},
    {{
      "type": "content",
      "section": "Section Name",
      "title": "Slide Title",
      "layout": "two_column",
      "left_title": "Left heading",
      "left_bullets": ["L1", "L2"],
      "right_title": "Right heading",
      "right_bullets": ["R1", "R2"]
    }},
    {{
      "type": "content",
      "section": "Section Name",
      "title": "Slide Title",
      "layout": "quote",
      "quote": "Impactful quote text here.",
      "attribution": "— Source"
    }},
    {{ "type": "closing", "title": "Thank you", "message": "Brief closing message" }}
  ]
}}

Rules:
- First and last slides are added automatically (title slide and closing). Do NOT include them in "slides".
- Include 1–3 section slides to divide content thematically.
- Mix layouts: use bullets for most content, two_column for comparisons, quote for 1-2 impactful moments.
- Target exactly {n_slides} content slides (excluding the auto title/closing).
- Choose theme colors appropriate for the topic. primary should be dark, accent should be vivid.
- Respond ONLY with valid JSON — no markdown fences, no commentary.\
"""


def generate_content(client: anthropic.Anthropic, topic: str, n_slides: int) -> dict:
    response = client.messages.create(
        model="claude-opus-4-7",
        max_tokens=4096,
        system=[{
            "type": "text",
            "text": SYSTEM_PROMPT.format(n_slides=n_slides),
            "cache_control": {"type": "ephemeral"},
        }],
        messages=[{
            "role": "user",
            "content": f'Create a {n_slides}-slide presentation about: "{topic}"',
        }],
    )
    raw = response.content[0].text.strip()
    if raw.startswith("```"):
        lines = raw.split("\n")
        raw = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
    return json.loads(raw)


def hex_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class SlideBuilder:
    def __init__(self, data: dict):
        self.data = data
        theme = data["theme"]
        self.primary = theme["primary"]
        self.accent = theme["accent"]
        self.tod = theme["text_on_dark"]
        self.tol = theme["text_on_light"]
        self.sections = [s["title"] for s in data["slides"] if s.get("type") == "section"]
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    def _blank(self):
        blank_layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(blank_layout)

    def _fill(self, slide, color: str):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_rgb(color)

    def _rect(self, slide, x, y, w, h, color: str):
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_rgb(color)
        shape.line.fill.background()
        return shape

    def _text(self, slide, text: str, x, y, w, h,
              size: int, bold: bool, color: str,
              align=PP_ALIGN.LEFT, wrap: bool = True):
        txb = slide.shapes.add_textbox(x, y, w, h)
        txb.word_wrap = wrap
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = hex_rgb(color)
        return txb

    # ── Title slide ──────────────────────────────────────────────────────────
    def add_title_slide(self):
        slide = self._blank()
        self._fill(slide, self.primary)
        self._rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, self.accent)
        self._rect(slide, Inches(0), Inches(5.6), SLIDE_W, Inches(0.06), self.accent)

        self._text(slide, self.data["title"],
                   Inches(0.5), Inches(1.8), Inches(12), Inches(2.2),
                   48, True, self.tod, PP_ALIGN.LEFT)
        self._text(slide, self.data.get("subtitle", ""),
                   Inches(0.5), Inches(4.0), Inches(10), Inches(0.9),
                   22, False, self.accent, PP_ALIGN.LEFT)

    # ── TOC slide ─────────────────────────────────────────────────────────────
    def add_toc_slide(self):
        slide = self._blank()
        self._fill(slide, self.primary)
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.3), self.accent)
        self._text(slide, "목 차", Inches(0.5), Inches(0.3), Inches(8), Inches(0.9),
                   28, True, self.primary, PP_ALIGN.LEFT)

        for i, sec in enumerate(self.sections):
            y = Inches(1.6) + i * Inches(0.82)
            self._rect(slide, Inches(0.5), y + Inches(0.15),
                       Inches(0.5), Inches(0.5), self.accent)
            self._text(slide, str(i + 1), Inches(0.5), y + Inches(0.1),
                       Inches(0.5), Inches(0.6), 18, True, self.primary, PP_ALIGN.CENTER)
            self._text(slide, sec, Inches(1.2), y, Inches(10), Inches(0.8),
                       20, False, self.tod, PP_ALIGN.LEFT)

    # ── Section divider ───────────────────────────────────────────────────────
    def add_section_slide(self, slide_data: dict):
        slide = self._blank()
        self._fill(slide, self.primary)
        self._rect(slide, Inches(0), Inches(0), Inches(6), SLIDE_H, self.accent)

        idx = self.sections.index(slide_data["title"]) + 1 if slide_data["title"] in self.sections else ""
        if idx:
            self._text(slide, f"{idx:02d}", Inches(0.3), Inches(0.3), Inches(5), Inches(1.5),
                       72, True, self.primary, PP_ALIGN.LEFT)
        self._text(slide, slide_data["title"],
                   Inches(0.3), Inches(2.5), Inches(5.5), Inches(2.5),
                   34, True, self.primary, PP_ALIGN.LEFT)
        self._text(slide, "─────────────────",
                   Inches(6.5), Inches(3.5), Inches(6), Inches(0.4),
                   14, False, self.tod, PP_ALIGN.LEFT)

    # ── Bullets layout ────────────────────────────────────────────────────────
    def _add_bullets_slide(self, slide_data: dict):
        slide = self._blank()
        self._fill(slide, "#FFFFFF")
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.4), self.primary)
        self._rect(slide, Inches(0), Inches(1.4), Inches(0.08), SLIDE_H - Inches(1.4), self.accent)

        self._text(slide, slide_data["title"],
                   Inches(0.3), Inches(0.2), Inches(12), Inches(1.0),
                   30, True, self.tod, PP_ALIGN.LEFT)

        bullets = slide_data.get("bullets", [])
        sub = slide_data.get("sub_bullets", {})
        y = Inches(1.7)
        for bullet in bullets:
            self._text(slide, f"▸  {bullet}",
                       Inches(0.4), y, Inches(12.4), Inches(0.55),
                       18, False, self.tol, PP_ALIGN.LEFT)
            y += Inches(0.58)
            for sub_b in sub.get(bullet, []):
                self._text(slide, f"      –  {sub_b}",
                           Inches(0.4), y, Inches(12.4), Inches(0.45),
                           15, False, "#555555", PP_ALIGN.LEFT)
                y += Inches(0.46)

        highlight = slide_data.get("highlight", "")
        if highlight:
            self._rect(slide, Inches(0.4), Inches(6.5), Inches(12.4), Inches(0.75), self.accent)
            self._text(slide, f"★  {highlight}",
                       Inches(0.6), Inches(6.55), Inches(12), Inches(0.65),
                       16, True, self.primary, PP_ALIGN.LEFT)

    # ── Two-column layout ─────────────────────────────────────────────────────
    def _add_two_column_slide(self, slide_data: dict):
        slide = self._blank()
        self._fill(slide, "#FFFFFF")
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.4), self.primary)
        self._rect(slide, Inches(6.5), Inches(1.4), Inches(0.06), SLIDE_H - Inches(1.4), "#DDDDDD")

        self._text(slide, slide_data["title"],
                   Inches(0.3), Inches(0.2), Inches(12), Inches(1.0),
                   30, True, self.tod, PP_ALIGN.LEFT)

        for col, cx in [("left", Inches(0.3)), ("right", Inches(6.8))]:
            heading = slide_data.get(f"{col}_title", "")
            items = slide_data.get(f"{col}_bullets", [])
            self._rect(slide, cx, Inches(1.5), Inches(5.9), Inches(0.5), self.accent)
            self._text(slide, heading, cx + Inches(0.1), Inches(1.55),
                       Inches(5.7), Inches(0.45), 18, True, self.primary, PP_ALIGN.LEFT)
            y = Inches(2.2)
            for item in items:
                self._text(slide, f"▸  {item}", cx, y, Inches(6.0), Inches(0.55),
                           16, False, self.tol, PP_ALIGN.LEFT)
                y += Inches(0.58)

    # ── Quote layout ──────────────────────────────────────────────────────────
    def _add_quote_slide(self, slide_data: dict):
        slide = self._blank()
        self._fill(slide, self.primary)
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.4), self.accent)

        self._text(slide, slide_data["title"],
                   Inches(0.3), Inches(0.2), Inches(12), Inches(1.0),
                   30, True, self.primary, PP_ALIGN.LEFT)
        self._text(slide, "“",
                   Inches(0.3), Inches(1.2), Inches(2), Inches(2),
                   96, True, self.accent, PP_ALIGN.LEFT)
        self._text(slide, slide_data.get("quote", ""),
                   Inches(1.2), Inches(2.5), Inches(10.8), Inches(3.0),
                   26, False, self.tod, PP_ALIGN.LEFT)
        self._text(slide, slide_data.get("attribution", ""),
                   Inches(8.0), Inches(6.2), Inches(5), Inches(0.7),
                   18, True, self.accent, PP_ALIGN.RIGHT)

    def add_content_slide(self, slide_data: dict):
        layout = slide_data.get("layout", "bullets")
        if layout == "two_column":
            self._add_two_column_slide(slide_data)
        elif layout == "quote":
            self._add_quote_slide(slide_data)
        else:
            self._add_bullets_slide(slide_data)

    # ── Closing slide ─────────────────────────────────────────────────────────
    def add_closing_slide(self, slide_data: dict):
        slide = self._blank()
        self._fill(slide, self.primary)
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.2), self.accent)
        self._rect(slide, Inches(0), SLIDE_H - Inches(0.2), SLIDE_W, Inches(0.2), self.accent)

        self._text(slide, slide_data.get("title", "감사합니다"),
                   Inches(0), Inches(2.3), SLIDE_W, Inches(2.0),
                   54, True, self.tod, PP_ALIGN.CENTER)
        self._text(slide, slide_data.get("message", ""),
                   Inches(1), Inches(4.6), Inches(11.3), Inches(1.2),
                   20, False, self.accent, PP_ALIGN.CENTER)

    # ── Assemble & save ───────────────────────────────────────────────────────
    def build(self, output_path: str) -> int:
        self.add_title_slide()
        if self.sections:
            self.add_toc_slide()

        closing_data = None
        for slide_data in self.data["slides"]:
            t = slide_data.get("type")
            if t == "section":
                self.add_section_slide(slide_data)
            elif t == "content":
                self.add_content_slide(slide_data)
            elif t == "closing":
                closing_data = slide_data

        self.add_closing_slide(closing_data or {"title": "감사합니다", "message": ""})

        self.prs.save(output_path)
        return len(self.prs.slides)


def main():
    parser = argparse.ArgumentParser(description="Generate a .pptx from a topic using Claude")
    parser.add_argument("topic", help="Presentation topic")
    parser.add_argument("--slides", type=int, default=8, metavar="N",
                        help="Number of content slides (default: 8)")
    parser.add_argument("--output", default="output.pptx", metavar="FILE",
                        help="Output file path (default: output.pptx)")
    parser.add_argument("--api-key", default="", metavar="KEY",
                        help="Anthropic API key (overrides ANTHROPIC_API_KEY env var)")
    args = parser.parse_args()

    api_key = args.api_key or os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("Error: Anthropic API key required.\n"
              "  Set ANTHROPIC_API_KEY env var, or pass --api-key <key>", file=sys.stderr)
        sys.exit(1)
    client = anthropic.Anthropic(api_key=api_key)

    print(f"Generating content for: {args.topic!r} ({args.slides} slides)...")
    data = generate_content(client, args.topic, args.slides)

    print("Building presentation...")
    builder = SlideBuilder(data)
    count = builder.build(args.output)

    print(f"Done! {count} slides saved to {args.output}")


if __name__ == "__main__":
    main()
